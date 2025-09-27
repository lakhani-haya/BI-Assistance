"""
Dashboard Export System
Provides export capabilities for dashboards in multiple formats (PDF, PNG, HTML, PowerPoint)
"""

import streamlit as st
import plotly.graph_objects as go
import pandas as pd
import json
import base64
import io
import zipfile
from datetime import datetime
import os
import tempfile
from typing import Dict, List, Any, Optional, Tuple, Union
from dataclasses import asdict
import uuid

# PDF generation
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# PowerPoint generation  
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Import dashboard components
from src.dashboard_builder import DashboardConfig, ChartConfig
from src.chart_editor import ChartStyling


class ExportFormat:
    """Export format options"""
    HTML = "html"
    PDF = "pdf"
    PNG = "png"
    SVG = "svg"
    POWERPOINT = "pptx"
    JSON = "json"
    EXCEL = "xlsx"


class DashboardExporter:
    """Dashboard export system with multiple format support"""
    
    def __init__(self, data: pd.DataFrame):
        self.data = data
        
    def render_export_interface(self, dashboard_config: DashboardConfig, 
                              chart_stylings: Dict[str, ChartStyling] = None) -> None:
        """Render export interface"""
        st.markdown("## 📤 Export Dashboard")
        
        if chart_stylings is None:
            chart_stylings = {}
        
        # Export format selection
        export_tabs = st.tabs([
            "📄 PDF Report",
            "🌐 Web Export", 
            "PowerPoint",
            "📱 Images",
            "💾 Data Export"
        ])
        
        with export_tabs[0]:
            self._render_pdf_export(dashboard_config, chart_stylings)
        
        with export_tabs[1]:
            self._render_web_export(dashboard_config, chart_stylings)
        
        with export_tabs[2]:
            self._render_powerpoint_export(dashboard_config, chart_stylings)
        
        with export_tabs[3]:
            self._render_image_export(dashboard_config, chart_stylings)
        
        with export_tabs[4]:
            self._render_data_export(dashboard_config)
    
    def _render_pdf_export(self, dashboard_config: DashboardConfig, 
                          chart_stylings: Dict[str, ChartStyling]) -> None:
        """Render PDF export interface"""
        st.markdown("### 📄 PDF Report Export")
        
        if not PDF_AVAILABLE:
            st.error("📋 PDF export requires `reportlab` package. Install with: `pip install reportlab`")
            return
        
        col1, col2 = st.columns(2)
        
        with col1:
            report_title = st.text_input("Report Title", value=dashboard_config.title)
            include_summary = st.checkbox("Include Executive Summary", value=True)
            include_data_info = st.checkbox("Include Data Information", value=True)
            
        with col2:
            page_size = st.selectbox("Page Size", ["A4", "Letter"], index=0)
            chart_size = st.selectbox("Chart Size", ["Small", "Medium", "Large"], index=1)
            include_insights = st.checkbox("Include AI Insights", value=True)
        
        # Report structure
        st.markdown("#### 📋 Report Structure")
        report_sections = st.multiselect(
            "Select Sections",
            [
                "Title Page",
                "Executive Summary", 
                "Data Overview",
                "Charts & Visualizations",
                "Key Insights",
                "Appendix"
            ],
            default=["Title Page", "Charts & Visualizations", "Key Insights"]
        )
        
        if st.button("📄 Generate PDF Report", type="primary"):
            with st.spinner("🔄 Generating PDF report..."):
                try:
                    pdf_bytes = self._generate_pdf_report(
                        dashboard_config, 
                        chart_stylings,
                        report_title,
                        page_size,
                        chart_size,
                        report_sections,
                        include_summary,
                        include_data_info,
                        include_insights
                    )
                    
                    if pdf_bytes:
                        st.download_button(
                            "📥 Download PDF Report",
                            data=pdf_bytes,
                            file_name=f"{dashboard_config.title}_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                            mime="application/pdf"
                        )
                        st.success("✅ PDF report generated successfully!")
                    else:
                        st.error("❌ Failed to generate PDF report")
                        
                except Exception as e:
                    st.error(f"❌ PDF generation error: {str(e)}")
    
    def _render_web_export(self, dashboard_config: DashboardConfig,
                          chart_stylings: Dict[str, ChartStyling]) -> None:
        """Render web export interface"""
        st.markdown("### 🌐 Web Export")
        
        col1, col2 = st.columns(2)
        
        with col1:
            include_interactive = st.checkbox("Include Interactive Features", value=True)
            include_filters = st.checkbox("Include Filters", value=True)
            responsive_design = st.checkbox("Responsive Design", value=True)
            
        with col2:
            theme_style = st.selectbox("Web Theme", ["Light", "Dark", "Auto"], index=0)
            include_css = st.checkbox("Include Custom CSS", value=False)
            standalone_html = st.checkbox("Standalone HTML", value=True)
        
        export_options = st.multiselect(
            "Export Options",
            ["HTML Dashboard", "JavaScript Code", "CSS Styles", "Complete Package"],
            default=["HTML Dashboard"]
        )
        
        if st.button("🌐 Generate Web Export", type="primary"):
            with st.spinner("🔄 Generating web export..."):
                try:
                    if "HTML Dashboard" in export_options:
                        html_content = self._generate_html_dashboard(
                            dashboard_config,
                            chart_stylings,
                            include_interactive,
                            include_filters,
                            responsive_design,
                            theme_style,
                            standalone_html
                        )
                        
                        st.download_button(
                            "📥 Download HTML Dashboard",
                            data=html_content,
                            file_name=f"{dashboard_config.title}_dashboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html",
                            mime="text/html"
                        )
                    
                    if "Complete Package" in export_options:
                        package_zip = self._generate_web_package(
                            dashboard_config,
                            chart_stylings,
                            include_interactive,
                            responsive_design
                        )
                        
                        st.download_button(
                            "📦 Download Complete Package",
                            data=package_zip,
                            file_name=f"{dashboard_config.title}_web_package_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                            mime="application/zip"
                        )
                    
                    st.success("✅ Web export generated successfully!")
                    
                except Exception as e:
                    st.error(f"❌ Web export error: {str(e)}")
    
    def _render_powerpoint_export(self, dashboard_config: DashboardConfig,
                                 chart_stylings: Dict[str, ChartStyling]) -> None:
        """Render PowerPoint export interface"""
        st.markdown("### 📊 PowerPoint Export")
        
        if not PPTX_AVAILABLE:
            st.error("📋 PowerPoint export requires `python-pptx` package. Install with: `pip install python-pptx`")
            return
        
        col1, col2 = st.columns(2)
        
        with col1:
            template_style = st.selectbox(
                "Presentation Template", 
                ["Corporate", "Modern", "Minimal", "Creative"],
                index=0
            )
            slides_per_chart = st.selectbox("Slides per Chart", [1, 2, 4], index=0)
            
        with col2:
            include_title_slide = st.checkbox("Include Title Slide", value=True)
            include_summary_slide = st.checkbox("Include Summary Slide", value=True)
            include_data_slides = st.checkbox("Include Data Tables", value=False)
        
        # Slide layout options
        st.markdown("#### 🎨 Slide Layout")
        slide_layout = st.selectbox(
            "Chart Layout",
            ["Full Slide", "Chart + Text", "2 Charts per Slide", "Chart + Data Table"],
            index=0
        )
        
        if st.button("📊 Generate PowerPoint", type="primary"):
            with st.spinner("🔄 Creating PowerPoint presentation..."):
                try:
                    pptx_bytes = self._generate_powerpoint_presentation(
                        dashboard_config,
                        chart_stylings,
                        template_style,
                        slide_layout,
                        slides_per_chart,
                        include_title_slide,
                        include_summary_slide,
                        include_data_slides
                    )
                    
                    if pptx_bytes:
                        st.download_button(
                            "📥 Download PowerPoint",
                            data=pptx_bytes,
                            file_name=f"{dashboard_config.title}_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        st.success("✅ PowerPoint presentation generated successfully!")
                    else:
                        st.error("❌ Failed to generate PowerPoint presentation")
                        
                except Exception as e:
                    st.error(f"❌ PowerPoint generation error: {str(e)}")
    
    def _render_image_export(self, dashboard_config: DashboardConfig,
                           chart_stylings: Dict[str, ChartStyling]) -> None:
        """Render image export interface"""
        st.markdown("### 📱 Image Export")
        
        col1, col2 = st.columns(2)
        
        with col1:
            image_format = st.selectbox("Image Format", ["PNG", "SVG", "JPEG"], index=0)
            image_quality = st.selectbox("Quality", ["High", "Medium", "Low"], index=0)
            
        with col2:
            image_size = st.selectbox(
                "Image Size",
                ["Small (800x600)", "Medium (1200x900)", "Large (1920x1080)", "Custom"],
                index=1
            )
            
            if image_size == "Custom":
                col_a, col_b = st.columns(2)
                with col_a:
                    custom_width = st.number_input("Width", min_value=400, max_value=4000, value=1200)
                with col_b:
                    custom_height = st.number_input("Height", min_value=300, max_value=3000, value=900)
        
        # Export options
        export_individual = st.checkbox("Export Individual Charts", value=True)
        export_combined = st.checkbox("Export Combined Dashboard", value=True)
        include_transparent = st.checkbox("Transparent Background", value=False)
        
        if st.button("📱 Generate Images", type="primary"):
            with st.spinner("🔄 Generating images..."):
                try:
                    # Parse image size
                    if image_size == "Custom":
                        width, height = custom_width, custom_height
                    else:
                        size_map = {
                            "Small (800x600)": (800, 600),
                            "Medium (1200x900)": (1200, 900),
                            "Large (1920x1080)": (1920, 1080)
                        }
                        width, height = size_map[image_size]
                    
                    # Quality mapping
                    quality_map = {"High": 300, "Medium": 150, "Low": 96}
                    dpi = quality_map[image_quality]
                    
                    images_zip = self._generate_image_export(
                        dashboard_config,
                        chart_stylings,
                        image_format.lower(),
                        width,
                        height,
                        dpi,
                        export_individual,
                        export_combined,
                        include_transparent
                    )
                    
                    if images_zip:
                        st.download_button(
                            "📥 Download Images",
                            data=images_zip,
                            file_name=f"{dashboard_config.title}_images_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                            mime="application/zip"
                        )
                        st.success("✅ Images generated successfully!")
                    else:
                        st.error("❌ Failed to generate images")
                        
                except Exception as e:
                    st.error(f"❌ Image generation error: {str(e)}")
    
    def _render_data_export(self, dashboard_config: DashboardConfig) -> None:
        """Render data export interface"""
        st.markdown("### 💾 Data Export")
        
        col1, col2 = st.columns(2)
        
        with col1:
            export_format = st.selectbox(
                "Export Format",
                ["Excel (XLSX)", "CSV", "JSON", "Parquet"],
                index=0
            )
            
            include_raw_data = st.checkbox("Include Raw Data", value=True)
            include_processed_data = st.checkbox("Include Processed Data", value=True)
            
        with col2:
            include_config = st.checkbox("Include Dashboard Config", value=True)
            include_metadata = st.checkbox("Include Metadata", value=True)
            compress_output = st.checkbox("Compress Output", value=False)
        
        if st.button("💾 Export Data", type="primary"):
            with st.spinner("🔄 Exporting data..."):
                try:
                    data_export = self._generate_data_export(
                        dashboard_config,
                        export_format,
                        include_raw_data,
                        include_processed_data,
                        include_config,
                        include_metadata,
                        compress_output
                    )
                    
                    if data_export:
                        file_extension = {
                            "Excel (XLSX)": "xlsx",
                            "CSV": "csv", 
                            "JSON": "json",
                            "Parquet": "parquet"
                        }[export_format]
                        
                        if compress_output:
                            file_extension = "zip"
                        
                        st.download_button(
                            "📥 Download Data Export",
                            data=data_export,
                            file_name=f"{dashboard_config.title}_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{file_extension}",
                            mime=self._get_mime_type(file_extension)
                        )
                        st.success("✅ Data exported successfully!")
                    else:
                        st.error("❌ Failed to export data")
                        
                except Exception as e:
                    st.error(f"❌ Data export error: {str(e)}")
    
    def _generate_pdf_report(self, dashboard_config: DashboardConfig,
                           chart_stylings: Dict[str, ChartStyling],
                           report_title: str, page_size: str, chart_size: str,
                           sections: List[str], include_summary: bool,
                           include_data_info: bool, include_insights: bool) -> bytes:
        """Generate PDF report"""
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                # Set page size
                pagesize = A4 if page_size == "A4" else letter
                doc = SimpleDocTemplate(tmp_file.name, pagesize=pagesize)
                
                # Create story (content)
                story = []
                styles = getSampleStyleSheet()
                
                # Custom styles
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Title'],
                    fontSize=24,
                    spaceAfter=30,
                    alignment=1  # Center
                )
                
                heading_style = ParagraphStyle(
                    'CustomHeading',
                    parent=styles['Heading1'],
                    fontSize=16,
                    spaceBefore=20,
                    spaceAfter=12
                )
                
                # Title page
                if "Title Page" in sections:
                    story.append(Paragraph(report_title, title_style))
                    story.append(Spacer(1, 0.5*inch))
                    story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
                    story.append(Spacer(1, 0.3*inch))
                    story.append(Paragraph(dashboard_config.description, styles['Normal']))
                    story.append(PageBreak())
                
                # Executive summary
                if "Executive Summary" in sections and include_summary:
                    story.append(Paragraph("Executive Summary", heading_style))
                    summary_text = self._generate_executive_summary(dashboard_config)
                    story.append(Paragraph(summary_text, styles['Normal']))
                    story.append(Spacer(1, 0.3*inch))
                
                # Data overview
                if "Data Overview" in sections and include_data_info:
                    story.append(Paragraph("Data Overview", heading_style))
                    data_info = self._generate_data_overview_text()
                    story.append(Paragraph(data_info, styles['Normal']))
                    story.append(Spacer(1, 0.3*inch))
                
                # Charts
                if "Charts & Visualizations" in sections:
                    story.append(Paragraph("Charts & Visualizations", heading_style))
                    
                    for chart_config in dashboard_config.charts:
                        # Add chart title
                        story.append(Paragraph(chart_config.title, styles['Heading2']))
                        
                        # Generate chart image
                        chart_image = self._generate_chart_image_for_pdf(
                            chart_config, 
                            chart_stylings.get(chart_config.chart_id),
                            chart_size
                        )
                        
                        if chart_image:
                            story.append(chart_image)
                        
                        story.append(Spacer(1, 0.2*inch))
                
                # Key insights
                if "Key Insights" in sections and include_insights:
                    story.append(Paragraph("Key Insights", heading_style))
                    insights_text = self._generate_insights_text(dashboard_config)
                    story.append(Paragraph(insights_text, styles['Normal']))
                
                # Build PDF
                doc.build(story)
                
                # Read the file
                with open(tmp_file.name, 'rb') as f:
                    pdf_bytes = f.read()
                
                # Clean up
                os.unlink(tmp_file.name)
                
                return pdf_bytes
                
        except Exception as e:
            st.error(f"PDF generation error: {str(e)}")
            return None
    
    def _generate_html_dashboard(self, dashboard_config: DashboardConfig,
                               chart_stylings: Dict[str, ChartStyling],
                               include_interactive: bool, include_filters: bool,
                               responsive_design: bool, theme_style: str,
                               standalone_html: bool) -> str:
        """Generate HTML dashboard"""
        # This is a simplified version - you would expand this significantly
        html_template = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{dashboard_config.title}</title>
    <script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
    <style>
        body {{
            font-family: Arial, sans-serif;
            margin: 0;
            padding: 20px;
            background-color: {'#f8f9fa' if theme_style == 'Light' else '#343a40' if theme_style == 'Dark' else '#ffffff'};
        }}
        .dashboard-header {{
            text-align: center;
            margin-bottom: 30px;
        }}
        .chart-container {{
            margin: 20px 0;
            padding: 15px;
            background: white;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }}
        .grid-container {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(400px, 1fr));
            gap: 20px;
        }}
    </style>
</head>
<body>
    <div class="dashboard-header">
        <h1>{dashboard_config.title}</h1>
        <p>{dashboard_config.description}</p>
    </div>
    
    <div class="grid-container">
"""
        
        # Add charts
        for i, chart_config in enumerate(dashboard_config.charts):
            html_template += f"""
        <div class="chart-container">
            <h3>{chart_config.title}</h3>
            <div id="chart_{i}"></div>
        </div>
"""
        
        html_template += """
    </div>
    
    <script>
        // Chart generation code would go here
        // This would include Plotly.js code to create interactive charts
    </script>
</body>
</html>
"""
        
        return html_template
    
    def _generate_powerpoint_presentation(self, dashboard_config: DashboardConfig,
                                        chart_stylings: Dict[str, ChartStyling],
                                        template_style: str, slide_layout: str,
                                        slides_per_chart: int, include_title_slide: bool,
                                        include_summary_slide: bool, 
                                        include_data_slides: bool) -> bytes:
        """Generate PowerPoint presentation"""
        try:
            prs = Presentation()
            
            # Title slide
            if include_title_slide:
                slide_layout_obj = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(slide_layout_obj)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = dashboard_config.title
                subtitle.text = f"{dashboard_config.description}\nGenerated on {datetime.now().strftime('%B %d, %Y')}"
            
            # Summary slide
            if include_summary_slide:
                slide_layout_obj = prs.slide_layouts[1]  # Content layout
                slide = prs.slides.add_slide(slide_layout_obj)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Dashboard Summary"
                summary_text = self._generate_executive_summary(dashboard_config)
                content.text = summary_text
            
            # Chart slides
            for chart_config in dashboard_config.charts:
                slide_layout_obj = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout_obj)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
                title_frame = title_shape.text_frame
                title_frame.text = chart_config.title
                title_frame.paragraphs[0].font.size = Pt(24)
                title_frame.paragraphs[0].font.bold = True
                
                # Add chart placeholder (in real implementation, you'd add actual chart image)
                chart_shape = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                chart_frame = chart_shape.text_frame
                chart_frame.text = f"Chart: {chart_config.chart_type.value.title()}\nData: {chart_config.x_column} vs {chart_config.y_column}"
            
            # Save to bytes
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                prs.save(tmp_file.name)
                
                with open(tmp_file.name, 'rb') as f:
                    pptx_bytes = f.read()
                
                os.unlink(tmp_file.name)
                
                return pptx_bytes
                
        except Exception as e:
            st.error(f"PowerPoint generation error: {str(e)}")
            return None
    
    def _generate_image_export(self, dashboard_config: DashboardConfig,
                             chart_stylings: Dict[str, ChartStyling],
                             image_format: str, width: int, height: int, dpi: int,
                             export_individual: bool, export_combined: bool,
                             include_transparent: bool) -> bytes:
        """Generate image export"""
        try:
            zip_buffer = io.BytesIO()
            
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                # Export individual charts
                if export_individual:
                    for i, chart_config in enumerate(dashboard_config.charts):
                        chart_image = self._create_chart_image(
                            chart_config,
                            chart_stylings.get(chart_config.chart_id),
                            image_format,
                            width,
                            height,
                            include_transparent
                        )
                        
                        if chart_image:
                            filename = f"chart_{i+1}_{chart_config.title.replace(' ', '_')}.{image_format}"
                            zip_file.writestr(filename, chart_image)
                
                # Export combined dashboard (placeholder)
                if export_combined:
                    # In real implementation, you'd create a combined dashboard image
                    combined_text = f"Combined dashboard image would be generated here"
                    zip_file.writestr("dashboard_combined.txt", combined_text)
            
            zip_buffer.seek(0)
            return zip_buffer.getvalue()
            
        except Exception as e:
            st.error(f"Image export error: {str(e)}")
            return None
    
    def _generate_data_export(self, dashboard_config: DashboardConfig,
                            export_format: str, include_raw_data: bool,
                            include_processed_data: bool, include_config: bool,
                            include_metadata: bool, compress_output: bool) -> bytes:
        """Generate data export"""
        try:
            if compress_output:
                zip_buffer = io.BytesIO()
                
                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                    # Raw data
                    if include_raw_data:
                        if export_format == "Excel (XLSX)":
                            excel_buffer = io.BytesIO()
                            self.data.to_excel(excel_buffer, index=False)
                            zip_file.writestr("raw_data.xlsx", excel_buffer.getvalue())
                        elif export_format == "CSV":
                            csv_data = self.data.to_csv(index=False)
                            zip_file.writestr("raw_data.csv", csv_data.encode())
                        elif export_format == "JSON":
                            json_data = self.data.to_json(orient='records', indent=2)
                            zip_file.writestr("raw_data.json", json_data.encode())
                    
                    # Dashboard config
                    if include_config:
                        config_data = asdict(dashboard_config)
                        config_json = json.dumps(config_data, indent=2, default=str)
                        zip_file.writestr("dashboard_config.json", config_json.encode())
                    
                    # Metadata
                    if include_metadata:
                        metadata = {
                            'export_date': datetime.now().isoformat(),
                            'data_shape': self.data.shape,
                            'data_columns': list(self.data.columns),
                            'export_format': export_format
                        }
                        metadata_json = json.dumps(metadata, indent=2)
                        zip_file.writestr("metadata.json", metadata_json.encode())
                
                zip_buffer.seek(0)
                return zip_buffer.getvalue()
            
            else:
                # Single file export
                if export_format == "Excel (XLSX)":
                    excel_buffer = io.BytesIO()
                    with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                        if include_raw_data:
                            self.data.to_excel(writer, sheet_name='Raw Data', index=False)
                        if include_config:
                            config_df = pd.DataFrame([asdict(dashboard_config)])
                            config_df.to_excel(writer, sheet_name='Config', index=False)
                    
                    excel_buffer.seek(0)
                    return excel_buffer.getvalue()
                
                elif export_format == "CSV":
                    return self.data.to_csv(index=False).encode()
                
                elif export_format == "JSON":
                    export_data = {}
                    if include_raw_data:
                        export_data['data'] = self.data.to_dict('records')
                    if include_config:
                        export_data['config'] = asdict(dashboard_config)
                    if include_metadata:
                        export_data['metadata'] = {
                            'export_date': datetime.now().isoformat(),
                            'data_shape': self.data.shape
                        }
                    
                    return json.dumps(export_data, indent=2, default=str).encode()
                
                return None
                
        except Exception as e:
            st.error(f"Data export error: {str(e)}")
            return None
    
    def _generate_executive_summary(self, dashboard_config: DashboardConfig) -> str:
        """Generate executive summary text"""
        return f"""
This dashboard presents {len(dashboard_config.charts)} key visualizations analyzing the dataset. 
The dashboard includes various chart types designed to provide comprehensive insights into the data patterns and trends.

Key highlights:
• {len(dashboard_config.charts)} interactive visualizations
• Multiple data perspectives and analysis angles
• Created on {dashboard_config.created_at.strftime('%B %d, %Y')}
• Theme: {dashboard_config.theme.value.title()}

The visualizations cover different aspects of the data to provide a complete analytical view.
"""
    
    def _generate_data_overview_text(self) -> str:
        """Generate data overview text"""
        return f"""
Dataset Overview:
• Total records: {len(self.data):,}
• Number of columns: {len(self.data.columns)}
• Data types: {', '.join(self.data.dtypes.value_counts().index.astype(str))}
• Memory usage: {self.data.memory_usage(deep=True).sum() / 1024 / 1024:.2f} MB

Column information:
{', '.join(self.data.columns[:10])}{'...' if len(self.data.columns) > 10 else ''}
"""
    
    def _generate_insights_text(self, dashboard_config: DashboardConfig) -> str:
        """Generate insights text"""
        return f"""
Key Insights from Analysis:

Based on the {len(dashboard_config.charts)} visualizations in this dashboard, several important patterns emerge from the data:

• The data contains {len(self.data)} records across {len(self.data.columns)} dimensions
• Multiple visualization types provide different analytical perspectives
• Interactive features enable detailed exploration of data patterns
• The dashboard design supports both high-level overview and detailed analysis

These insights provide a foundation for data-driven decision making and further analysis.
"""
    
    def _generate_chart_image_for_pdf(self, chart_config: ChartConfig,
                                    styling: ChartStyling, chart_size: str) -> Optional[Image]:
        """Generate chart image for PDF"""
        # Placeholder - in real implementation, you'd generate actual chart images
        return None
    
    def _create_chart_image(self, chart_config: ChartConfig, styling: ChartStyling,
                          image_format: str, width: int, height: int,
                          include_transparent: bool) -> bytes:
        """Create chart image"""
        # Placeholder - in real implementation, you'd generate actual chart images
        return b"chart_image_placeholder"
    
    def _generate_web_package(self, dashboard_config: DashboardConfig,
                            chart_stylings: Dict[str, ChartStyling],
                            include_interactive: bool, responsive_design: bool) -> bytes:
        """Generate complete web package"""
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            # Add HTML file
            html_content = self._generate_html_dashboard(
                dashboard_config, chart_stylings, include_interactive,
                True, responsive_design, "Light", True
            )
            zip_file.writestr("dashboard.html", html_content.encode())
            
            # Add CSS file
            css_content = self._generate_css_styles(dashboard_config, responsive_design)
            zip_file.writestr("styles.css", css_content.encode())
            
            # Add JavaScript file
            js_content = self._generate_javascript_code(dashboard_config, chart_stylings)
            zip_file.writestr("dashboard.js", js_content.encode())
            
            # Add README
            readme_content = self._generate_web_readme(dashboard_config)
            zip_file.writestr("README.md", readme_content.encode())
        
        zip_buffer.seek(0)
        return zip_buffer.getvalue()
    
    def _generate_css_styles(self, dashboard_config: DashboardConfig, responsive: bool) -> str:
        """Generate CSS styles"""
        return """
/* Dashboard Styles */
body {
    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    margin: 0;
    padding: 20px;
    background-color: #f8f9fa;
}

.dashboard-container {
    max-width: 1200px;
    margin: 0 auto;
}

.chart-container {
    background: white;
    border-radius: 8px;
    padding: 20px;
    margin: 20px 0;
    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
}

@media (max-width: 768px) {
    .dashboard-container {
        padding: 10px;
    }
    
    .chart-container {
        margin: 10px 0;
        padding: 15px;
    }
}
"""
    
    def _generate_javascript_code(self, dashboard_config: DashboardConfig,
                                chart_stylings: Dict[str, ChartStyling]) -> str:
        """Generate JavaScript code"""
        return """
// Dashboard JavaScript
document.addEventListener('DOMContentLoaded', function() {
    console.log('Dashboard loaded');
    
    // Initialize charts
    initializeCharts();
    
    // Add interactivity
    setupInteractivity();
});

function initializeCharts() {
    // Chart initialization code would go here
}

function setupInteractivity() {
    // Interactivity setup code would go here
}
"""
    
    def _generate_web_readme(self, dashboard_config: DashboardConfig) -> str:
        """Generate README for web package"""
        return f"""# {dashboard_config.title}

{dashboard_config.description}

## Files Included

- `dashboard.html` - Main dashboard file
- `styles.css` - Dashboard styles
- `dashboard.js` - Interactive functionality
- `README.md` - This file

## Usage

Open `dashboard.html` in a web browser to view the dashboard.

## Generated

Created on {datetime.now().strftime('%B %d, %Y')} using Smart BI Assistant.
"""
    
    def _get_mime_type(self, file_extension: str) -> str:
        """Get MIME type for file extension"""
        mime_types = {
            'pdf': 'application/pdf',
            'html': 'text/html',
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'csv': 'text/csv',
            'json': 'application/json',
            'zip': 'application/zip',
            'png': 'image/png',
            'svg': 'image/svg+xml',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        return mime_types.get(file_extension, 'application/octet-stream')


# Export main classes
__all__ = [
    'ExportFormat',
    'DashboardExporter'
]
